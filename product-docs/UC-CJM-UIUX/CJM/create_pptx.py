"""
Create PowerPoint presentation from customer-journey-map-presentation.md
Using python-pptx library
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os
from xml.scripts.saxutils import escape

def rgb_to_int(r, g, b):
    """Convert RGB to integer for pptx"""
    return (r << 16) | (g << 8) | b

# Colors as RGB tuples
COLOR_STAGE_BG = (232, 245, 233)  # Light green
COLOR_STAGE_BORDER = (42, 95, 60)  # Dark green
COLOR_ACTION_BG = (227, 242, 253)  # Light blue
COLOR_ACTION_BORDER = (25, 118, 210)  # Blue
COLOR_INTERVENTION_BG = (255, 243, 224)  # Light orange
COLOR_INTERVENTION_BORDER = (245, 124, 0)  # Orange
COLOR_TEXT = (50, 50, 50)
COLOR_TEXT_GRAY = (120, 120, 120)
COLOR_TEXT_LIGHT = (100, 100, 100)

def create_slide_with_diagram(prs):
    """Create a single slide with CJM diagram"""
    
    # Add blank slide
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Healora AI - Customer Journey Map"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLOR_STAGE_BORDER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(0.3))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "From Registration to Longevity Path"
    subtitle_frame.paragraphs[0].font.size = Pt(14)
    subtitle_frame.paragraphs[0].font.color.rgb = RgbColor(100, 100, 100)
    
    # Stage positions
    stage_y = 1.2
    stage_height = 3.5
    stage_width = 1.6
    stage_gap = 0.3
    start_x = 0.5
    
    stages = [
        {
            "title": "1. REGISTRATION",
            "color": COLOR_STAGE_BG,
            "border": COLOR_STAGE_BORDER,
            "items": ["Landing Page", "Sign Up", "Email Verify", "Account Active"],
            "sub_items": ["< 5 min", "Data: Email, Name"]
        },
        {
            "title": "2. ONBOARDING",
            "color": COLOR_STAGE_BG,
            "border": COLOR_STAGE_BORDER,
            "items": ["Health Profile", "Health Quiz", "Wearable Connect", "Blood Test", "Digital Twin Init"],
            "sub_items": ["24-48h Twin gen"]
        },
        {
            "title": "3. ENRICHING",
            "color": COLOR_STAGE_BG,
            "border": COLOR_STAGE_BORDER,
            "items": ["Daily Diary", "Weekly Check-in", "Wearable Sync", "Diet Track", "Twin Update"],
            "sub_items": ["2-4 weeks cycle"]
        },
        {
            "title": "4. RECOMMENDATIONS",
            "color": COLOR_STAGE_BG,
            "border": COLOR_STAGE_BORDER,
            "items": ["AI/ML Generation", "Short-term (1-4w)", "Long-term (3-12m)", "Review & Track"],
            "sub_items": ["Sleep, Nutrition, IF", "HIIT, Supplements"]
        },
        {
            "title": "5. OPTIMIZATION",
            "color": COLOR_STAGE_BG,
            "border": COLOR_STAGE_BORDER,
            "items": ["Week 1-4", "Week 5-8", "Week 9-12", "Month 3-6", "Month 6-12"],
            "sub_items": ["Longevity focus"]
        }
    ]
    
    # Draw stages
    for i, stage in enumerate(stages):
        x = start_x + i * (stage_width + stage_gap)
        
        # Stage container
        stage_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(stage_y), Inches(stage_width), Inches(stage_height)
        )
        stage_shape.fill.solid()
        stage_shape.fill.fore_color.rgb = stage["color"]
        stage_shape.line.color.rgb = stage["border"]
        stage_shape.line.width = Pt(2)
        
        # Stage title
        title_shape = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(stage_y + 0.1), 
            Inches(stage_width - 0.2), Inches(0.3)
        )
        tf = title_shape.text_frame
        tf.text = stage["title"]
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = stage["border"]
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Items
        item_y = stage_y + 0.5
        for item in stage["items"]:
            item_shape = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(item_y),
                Inches(stage_width - 0.2), Inches(0.25)
            )
            tf = item_shape.text_frame
            tf.text = "• " + item
            tf.paragraphs[0].font.size = Pt(8)
            tf.paragraphs[0].font.color.rgb = RgbColor(50, 50, 50)
            item_y += 0.28
        
        # Sub-items (additional info)
        sub_y = item_y + 0.1
        for sub in stage["sub_items"]:
            sub_shape = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(sub_y),
                Inches(stage_width - 0.2), Inches(0.2)
            )
            tf = sub_shape.text_frame
            tf.text = "(" + sub + ")"
            tf.paragraphs[0].font.size = Pt(7)
            tf.paragraphs[0].font.italic = True
            tf.paragraphs[0].font.color.rgb = RgbColor(120, 120, 120)
            sub_y += 0.18
        
        # Arrow between stages
        if i < len(stages) - 1:
            arrow_x = x + stage_width + 0.05
            arrow_shape = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(arrow_x), Inches(stage_y + stage_height/2 - 0.15),
                Inches(0.2), Inches(0.3)
            )
            arrow_shape.fill.solid()
            arrow_shape.fill.fore_color.rgb = COLOR_STAGE_BORDER
            arrow_shape.line.color.rgb = COLOR_STAGE_BORDER
    
    # Interventions boxes (below stages)
    int_y = 4.9
    int_height = 0.9
    
    # Short-term box
    short_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.8), Inches(int_y), Inches(2.5), Inches(int_height)
    )
    short_box.fill.solid()
    short_box.fill.fore_color.rgb = COLOR_INTERVENTION_BG
    short_box.line.color.rgb = COLOR_INTERVENTION_BORDER
    short_box.line.width = Pt(1)
    
    short_text = slide.shapes.add_textbox(
        Inches(2.9), Inches(int_y + 0.1),
        Inches(2.3), Inches(0.8)
    )
    tf = short_text.text_frame
    tf.text = "Short-term (1-4 weeks)\n• Sleep Protocol 22:00-23:00\n• Nutrition: No late eating\n• Movement: 10k steps\n• Meditation: 10min/day"
    tf.paragraphs[0].font.size = Pt(7)
    tf.paragraphs[0].font.bold = True
    
    # Long-term box
    long_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.5), Inches(int_y), Inches(2.8), Inches(int_height)
    )
    long_box.fill.solid()
    long_box.fill.fore_color.rgb = COLOR_INTERVENTION_BG
    long_box.line.color.rgb = COLOR_INTERVENTION_BORDER
    long_box.line.width = Pt(1)
    
    long_text = slide.shapes.add_textbox(
        Inches(5.6), Inches(int_y + 0.1),
        Inches(2.6), Inches(0.8)
    )
    tf = long_text.text_frame
    tf.text = "Long-term (3-12 months)\n• IF Protocol 16:8\n• Personalized Supplements\n• HIIT 2x/week\n• Quarterly Lab Review"
    tf.paragraphs[0].font.size = Pt(7)
    tf.paragraphs[0].font.bold = True
    
    # Metrics bar at bottom
    metrics_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(6.0), Inches(9), Inches(0.6)
    )
    metrics_box.fill.solid()
    metrics_box.fill.fore_color.rgb = RgbColor(236, 239, 241)
    metrics_box.line.color.rgb = RgbColor(69, 90, 100)
    
    metrics_text = slide.shapes.add_textbox(
        Inches(0.6), Inches(6.1),
        Inches(8.8), Inches(0.4)
    )
    tf = metrics_text.text_frame
    tf.text = "Target Metrics: Sign-up 30% | Onboarding Completion 70% | Diary Adherence 80% | Recommendation Acceptance 60% | 6-month Retention 50%"
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RgbColor(50, 50, 50)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Goal text
    goal_text = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.7),
        Inches(9), Inches(0.3)
    )
    tf = goal_text.text_frame
    tf.text = "GOAL: Obesity diagnosis → Personalized Longevity Path | VALUE: AI-powered Digital Twin for interventions"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_STAGE_BORDER
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Create the CJM slide
    create_slide_with_diagram(prs)
    
    # Save
    output_path = os.path.join(
        os.path.dirname(__file__),
        "..", "..", "product-docs", "product", "CJM", 
        "customer-journey-map.pptx"
    )
    output_path = os.path.abspath(output_path)
    
    prs.save(output_path)
    print(f"Saved: {output_path}")

if __name__ == "__main__":
    main()