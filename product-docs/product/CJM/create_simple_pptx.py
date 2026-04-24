"""
Create simple CJM PowerPoint slide
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Add blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Healora AI - Customer Journey Map"
p.font.size = Pt(28)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Subtitle
sub_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.7), Inches(9.4), Inches(0.4))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "From Registration to Longevity Path"
p.font.size = Pt(16)
p.alignment = PP_ALIGN.CENTER

# Stage data
stages = [
    ("1. REGISTRATION", ["Landing Page", "Sign Up", "Email Verify", "Account Active"], "< 5 min"),
    ("2. ONBOARDING", ["Health Profile", "Health Quiz", "Wearable Connect", "Blood Test", "Digital Twin Init"], "24-48h"),
    ("3. ENRICHING", ["Daily Diary", "Weekly Check-in", "Wearable Sync", "Diet Track", "Twin Update"], "2-4 weeks"),
    ("4. RECOMMENDATIONS", ["AI/ML Generation", "Review & Track", "Short-term (1-4w)", "Long-term (3-12m)"], "Ongoing"),
    ("5. OPTIMIZATION", ["Week 1-4", "Week 5-8", "Week 9-12", "Month 3-6", "Month 6-12"], "6-12 months"),
]

# Colors (RGBColor objects)
stage_colors = [
    (RGBColor(232, 245, 233), RGBColor(42, 95, 60)),   # Green
    (RGBColor(227, 242, 253), RGBColor(25, 118, 210)), # Blue
    (RGBColor(243, 229, 245), RGBColor(123, 31, 162)), # Purple
    (RGBColor(255, 243, 224), RGBColor(245, 124, 0)), # Orange
    (RGBColor(255, 235, 238), RGBColor(198, 40, 40)),  # Red
]

# Draw stages
x_start = 0.4
y_start = 1.3
box_width = 1.7
box_height = 3.8
gap = 0.25

for i, (title, items, duration) in enumerate(stages):
    x = x_start + i * (box_width + gap)
    bg_color, border_color = stage_colors[i]
    
    # Stage box
    stage = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y_start),
        Inches(box_width), Inches(box_height)
    )
    stage.fill.solid()
    stage.fill.fore_color.rgb = bg_color
    stage.line.color.rgb = border_color
    stage.line.width = Pt(2)
    
    # Title inside box
    title_tb = slide.shapes.add_textbox(
        Inches(x + 0.1), Inches(y_start + 0.1),
        Inches(box_width - 0.2), Inches(0.4)
    )
    tf = title_tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = border_color
    
    # Items
    item_y = y_start + 0.6
    for item in items:
        item_tb = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(item_y),
            Inches(box_width - 0.2), Inches(0.3)
        )
        tf = item_tb.text_frame
        p = tf.paragraphs[0]
        p.text = "* " + item
        p.font.size = Pt(9)
        item_y += 0.35
    
    # Duration
    dur_tb = slide.shapes.add_textbox(
        Inches(x + 0.1), Inches(y_start + box_height - 0.5),
        Inches(box_width - 0.2), Inches(0.3)
    )
    tf = dur_tb.text_frame
    p = tf.paragraphs[0]
    p.text = "(" + duration + ")"
    p.font.size = Pt(8)
    p.font.italic = True
    
    # Arrow to next
    if i < len(stages) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(x + box_width + 0.05), Inches(y_start + box_height/2 - 0.15),
            Inches(0.15), Inches(0.3)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)
        arrow.line.fill.background()

# Short-term interventions box
short_y = 5.3
short_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(2.5), Inches(short_y), Inches(2.2), Inches(0.9)
)
short_box.fill.solid()
short_box.fill.fore_color.rgb = RGBColor(255, 243, 224)
short_box.line.width = Pt(1)

short_title = slide.shapes.add_textbox(
    Inches(2.6), Inches(short_y + 0.05),
    Inches(2.0), Inches(0.3)
)
tf = short_title.text_frame
p = tf.paragraphs[0]
p.text = "Short-term (1-4 weeks)"
p.font.size = Pt(10)
p.font.bold = True

short_items = slide.shapes.add_textbox(
    Inches(2.6), Inches(short_y + 0.35),
    Inches(2.0), Inches(0.6)
)
tf = short_items.text_frame
p = tf.paragraphs[0]
p.text = "Sleep: 22:00-23:00"
p = tf.add_paragraph()
p.text = "Nutrition: No late eating"
p = tf.add_paragraph()
p.text = "Movement: 10k steps"
p = tf.add_paragraph()
p.text = "Meditation: 10min"
for para in tf.paragraphs:
    para.font.size = Pt(8)

# Long-term interventions box
long_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(5.0), Inches(short_y), Inches(2.5), Inches(0.9)
)
long_box.fill.solid()
long_box.fill.fore_color.rgb = RGBColor(255, 243, 224)
long_box.line.width = Pt(1)

long_title = slide.shapes.add_textbox(
    Inches(5.1), Inches(short_y + 0.05),
    Inches(2.3), Inches(0.3)
)
tf = long_title.text_frame
p = tf.paragraphs[0]
p.text = "Long-term (3-12 months)"
p.font.size = Pt(10)
p.font.bold = True

long_items = slide.shapes.add_textbox(
    Inches(5.1), Inches(short_y + 0.35),
    Inches(2.3), Inches(0.6)
)
tf = long_items.text_frame
p = tf.paragraphs[0]
p.text = "IF Protocol 16:8"
p = tf.add_paragraph()
p.text = "Supplements (personalized)"
p = tf.add_paragraph()
p.text = "HIIT 2x/week"
p = tf.add_paragraph()
p.text = "Quarterly Lab Review"
for para in tf.paragraphs:
    para.font.size = Pt(8)

# Metrics bar
metrics = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0.3), Inches(6.4), Inches(9.4), Inches(0.5)
)
metrics.fill.solid()
metrics.fill.fore_color.rgb = RGBColor(236, 239, 241)
metrics.line.fill.background()

metrics_text = slide.shapes.add_textbox(
    Inches(0.4), Inches(6.5), Inches(9.2), Inches(0.3)
)
tf = metrics_text.text_frame
p = tf.paragraphs[0]
p.text = "Target Metrics: Sign-up 30% | Onboarding 70% | Adherence 80% | Acceptance 60% | Retention 50%"
p.font.size = Pt(11)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Goal
goal = slide.shapes.add_textbox(
    Inches(0.3), Inches(7.0), Inches(9.4), Inches(0.3))
tf = goal.text_frame
p = tf.paragraphs[0]
p.text = "GOAL: Obesity diagnosis to Personalized Longevity Path | VALUE: AI Digital Twin"
p.font.size = Pt(12)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Save
output_path = os.path.join(os.path.dirname(__file__), "customer-journey-map.pptx")
prs.save(output_path)
print(f"Created: {output_path}")